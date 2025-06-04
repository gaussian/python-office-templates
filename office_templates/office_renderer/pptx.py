from pptx import Presentation
from .charts import process_chart
from .images import (
    replace_shape_with_image,
    should_replace_shape_with_image,
)
from .loops import (
    is_loop_end,
    is_loop_start,
    process_loops,
    LOOP_START_PATTERN,
    LOOP_END_PATTERN,
)
from .paragraphs import process_paragraph
from .pptx_utils import remove_shape
from .tables import process_table_cell


def clear_loop_directives(prs):
    """
    Clear the text of all shapes that contain loop directives.

    Args:
        prs: The Presentation object
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                text = shape.text_frame.text.strip()
                if LOOP_START_PATTERN.search(text) or LOOP_END_PATTERN.search(text):
                    # Clear text at paragraph level to handle formatting
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            for run in paragraph.runs:
                                run.text = ""
                        else:
                            paragraph.text = ""


def render_pptx(template, context: dict, output, perm_user):
    """
    Render the PPTX template (a path string or a file-like object) using the provided context and save to output.
    'output' can be a path string or a file-like object. If it's a file-like object, it will be rewound after saving.
    """
    # Support template as a file path or file-like object.
    if isinstance(template, str):
        prs = Presentation(template)
    else:
        template.seek(0)
        prs = Presentation(template)

    errors = []

    # Process loops first - identify loop sections and duplicate slides
    slides_to_process = process_loops(prs, context, perm_user, errors)

    # Process all slides including duplicated ones from loops
    for slide_info in slides_to_process:
        slide = slide_info["slide"]
        slide_number = slide_info.get("slide_number", 0)
        extra_context = slide_info.get("extra_context", {})

        # Create slide context (include `extra_context`, which is where loop variables are)
        slide_context = {
            **context,
            **extra_context,
            "slide_number": slide_number,
        }

        # Add loop variable to context if present
        if "loop_var" in slide_info and "loop_item" in slide_info:
            slide_context[slide_info["loop_var"]] = slide_info["loop_item"]

        # Process the slide
        process_single_slide(slide, slide_context, slide_number, perm_user, errors)

    if errors:
        print("Rendering aborted due to the following errors:")
        for err in set(errors):
            print(f" - {err}")
        print("Output file not saved.")
        return None, errors

    # Save to output (file path or file-like object)
    if isinstance(output, str):
        prs.save(output)
    else:
        prs.save(output)
        output.seek(0)

    return output, None


def process_single_slide(slide, context, slide_number, perm_user, errors, placeholders=None):
    """Process a single slide with the given context and placeholders."""
    # Handle placeholders first if provided
    if placeholders:
        _process_placeholders(slide, placeholders, context, slide_number, perm_user, errors)
    
    # Process the slide's shapes
    for shape in slide.shapes:
        # Skip loop directive shapes - we'll clear them later
        if is_loop_start(shape) or is_loop_end(shape):
            continue

        # Process the shape content
        process_shape_content(
            shape, slide, context, slide_number, perm_user, errors
        )


def _process_placeholders(slide, placeholders, context, slide_number, perm_user, errors):
    """Process placeholder shapes on a slide."""
    from ..templating import process_text
    
    placeholder_index = 0
    for shape in slide.shapes:
        # Check if this shape is a placeholder
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder and placeholder_index < len(placeholders):
            try:
                placeholder_text = placeholders[placeholder_index]
                processed_text = process_text(placeholder_text, context, perm_user)
                
                # Set the processed text to the shape
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = processed_text
                elif hasattr(shape, 'text'):
                    shape.text = processed_text
                    
                placeholder_index += 1
            except Exception as e:
                errors.append(f"Error processing placeholder {placeholder_index} (slide {slide_number}): {e}")


def process_shape_content(shape, slide, context, slide_number, perm_user, errors):
    """Process the content of a shape based on its type."""
    # 1) Check if this shape should be replaced with an image.
    if should_replace_shape_with_image(shape):
        try:
            replace_shape_with_image(
                shape,
                slide,
                context=context,
                perm_user=perm_user,
            )
        except Exception as e:
            errors.append(f"Error processing image (slide {slide_number}): {e}")
        # Skip further processing for this shape.
        return

    # 2) Check if this shape should be removed (because it's a loop directive).
    if is_loop_start(shape) or is_loop_end(shape):
        remove_shape(shape)
        return

    # 3) Process text frames (non-table).
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            # Merge any placeholders that are split across multiple runs.
            try:
                process_paragraph(
                    paragraph=paragraph,
                    context=context,
                    perm_user=perm_user,
                    mode="normal",  # for text frames
                )
            except Exception as e:
                errors.append(f"Error in paragraph (slide {slide_number}): {e}")

    # 4) Process tables.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                try:
                    process_table_cell(
                        cell=cell,
                        context=context,
                        perm_user=perm_user,
                    )
                except Exception as e:
                    errors.append(f"Error in table cell (slide {slide_number}): {e}")

    # 5) Process chart spreadsheets.
    if getattr(shape, "has_chart", False):
        try:
            process_chart(
                chart=shape.chart,
                context=context,
                perm_user=perm_user,
            )
        except Exception as e:
            errors.append(f"Error in chart (slide {slide_number}): {e}")


def build_layout_mapping(template_files, use_tagged_layouts=False, use_all_slides_as_layouts_by_title=False):
    """
    Build a mapping of layout IDs to slide objects from multiple template files.
    
    Args:
        template_files: List of template file paths or file-like objects
        use_tagged_layouts: If True, include slides with % layout XXX % tags
        use_all_slides_as_layouts_by_title: If True, use all slides as layouts by title
        
    Returns:
        dict: Mapping of layout ID (str) to (presentation, slide) tuple
    """
    layout_mapping = {}
    
    for template_file in template_files:
        # Load presentation
        if isinstance(template_file, str):
            prs = Presentation(template_file)
        else:
            template_file.seek(0)
            prs = Presentation(template_file)
        
        # Get master layouts
        master_layouts = get_master_layouts(prs)
        for layout_id, layout in master_layouts.items():
            layout_mapping[layout_id] = (prs, layout)
        
        # Get tagged layouts if enabled
        if use_tagged_layouts:
            tagged_layouts = get_tagged_layouts(prs)
            for layout_id, slide in tagged_layouts.items():
                layout_mapping[layout_id] = (prs, slide)
        
        # Get title layouts if enabled
        if use_all_slides_as_layouts_by_title:
            title_layouts = get_title_layouts(prs)
            for layout_id, slide in title_layouts.items():
                layout_mapping[layout_id] = (prs, slide)
    
    return layout_mapping


def get_master_layouts(prs):
    """Get master layout slides where ID is the layout name."""
    layouts = {}
    for slide_layout in prs.slide_layouts:
        # Use the layout name as the ID
        layout_id = slide_layout.name
        # For master layouts, we'll store the layout itself
        # and create slides from it when needed
        layouts[layout_id] = slide_layout
    
    return layouts


def get_tagged_layouts(prs):
    """Get slides that have shapes with % layout XXX % tags."""
    import re
    layouts = {}
    pattern = re.compile(r'%\s*layout\s+(\w+)\s*%', re.IGNORECASE)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                text = shape.text_frame.text.strip()
                match = pattern.search(text)
                if match:
                    layout_id = match.group(1)
                    layouts[layout_id] = slide
                    break  # Found the layout tag, move to next slide
    
    return layouts


def get_title_layouts(prs):
    """Get all slides as layouts where ID is the slide title."""
    layouts = {}
    
    for slide in prs.slides:
        # Find the title shape (usually the first shape or a shape with specific placeholder type)
        title = None
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                # Use the first text shape as title
                title = shape.text_frame.text.strip()
                break
        
        if title:
            layouts[title] = slide
    
    return layouts


def compose_pptx(template_files, slides, global_context, output, perm_user=None, 
                 use_tagged_layouts=False, use_all_slides_as_layouts_by_title=False):
    """
    Compose a PPTX deck using layout slides from multiple template files.
    
    Args:
        template_files: List of template file paths or file-like objects
        slides: List of slide dictionaries, each containing 'layout' key and context data
        global_context: Global context dictionary
        output: Output file path or file-like object
        perm_user: Permission user object
        use_tagged_layouts: If True, include slides with % layout XXX % tags
        use_all_slides_as_layouts_by_title: If True, use all slides as layouts by title
        
    Returns:
        Tuple of (output, errors) - errors is None if successful, list of errors otherwise
    """
    errors = []
    
    try:
        # Validate inputs
        if not template_files:
            errors.append("No template files provided")
            return None, errors
        
        if not slides:
            errors.append("No slides specified")
            return None, errors
        
        # Build layout mapping from all template files
        layout_mapping = build_layout_mapping(
            template_files, 
            use_tagged_layouts=use_tagged_layouts,
            use_all_slides_as_layouts_by_title=use_all_slides_as_layouts_by_title
        )
        
        if not layout_mapping:
            errors.append("No layout slides found in template files")
            return None, errors
        
        # Use the first template as the base presentation
        if isinstance(template_files[0], str):
            base_prs = Presentation(template_files[0])
        else:
            template_files[0].seek(0)
            base_prs = Presentation(template_files[0])
        
        # Remove all slides from the base presentation to create a blank deck
        sld_ids = base_prs.slides._sldIdLst
        while len(sld_ids) > 0:
            slide_id = sld_ids[0]
            sld_ids.remove(slide_id)
        
        # Create slides from the slide specifications
        for slide_index, slide_spec in enumerate(slides):
            try:
                # Validate slide specification
                if 'layout' not in slide_spec:
                    errors.append(f"Slide {slide_index + 1}: Missing 'layout' key")
                    continue
                
                layout_id = slide_spec['layout']
                if layout_id not in layout_mapping:
                    errors.append(f"Slide {slide_index + 1}: Layout '{layout_id}' not found")
                    continue
                
                # Get the layout slide or layout
                source_prs, layout_item = layout_mapping[layout_id]
                
                # Handle both slide layouts and actual slides
                if hasattr(layout_item, 'slide_layout'):
                    # It's a slide - copy it directly
                    from .pptx_utils import copy_slide_across_presentations
                    new_slide = copy_slide_across_presentations(base_prs, layout_item)
                else:
                    # It's a slide layout - find equivalent layout in base presentation or use blank
                    if len(base_prs.slide_layouts) > 6:
                        layout = base_prs.slide_layouts[6]  # Use blank layout
                    elif len(base_prs.slide_layouts) > 0:
                        layout = base_prs.slide_layouts[0]  # Use first available layout
                    else:
                        errors.append(f"Slide {slide_index + 1}: No slide layouts available in base presentation")
                        continue
                    new_slide = base_prs.slides.add_slide(layout)
                
                # Prepare slide context
                slide_context = {**global_context, **slide_spec}
                slide_number = slide_index + 1
                
                # Get placeholders if specified
                placeholders = slide_spec.get('placeholders', None)
                
                # Process the slide
                process_single_slide(new_slide, slide_context, slide_number, perm_user, errors, placeholders)
                
            except Exception as e:
                errors.append(f"Error processing slide {slide_index + 1}: {e}")
        
        # If we have errors, don't save
        if errors:
            print("Composition aborted due to the following errors:")
            for err in set(errors):
                print(f" - {err}")
            print("Output file not saved.")
            return None, errors
        
        # Save to output (file path or file-like object)
        if isinstance(output, str):
            base_prs.save(output)
        else:
            base_prs.save(output)
            output.seek(0)
        
        return output, None
        
    except Exception as e:
        errors.append(f"Unexpected error during composition: {e}")
        return None, errors
