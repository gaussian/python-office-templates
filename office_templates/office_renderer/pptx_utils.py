"""
Utility functions for PPTX manipulation.
"""

from copy import deepcopy
from pptx import Presentation


def duplicate_slide(pres, index):
    """
    Create a duplicate of the slide at the given index and return the new slide.
    
    This function duplicates a slide by:
    1. Getting the source slide
    2. Creating a new blank slide
    3. Copying all shapes from source to destination
    
    Args:
        pres: The Presentation object
        index: Index of the slide to duplicate
        
    Returns:
        The newly created slide object
    """
    source = pres.slides[index]
    blank_slide_layout = pres.slide_layouts[6]  # typically a blank layout
    new_slide = pres.slides.add_slide(blank_slide_layout)

    for shape in source.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def remove_shape(slide, shape):
    """
    Remove a shape from a slide.
    
    This function removes a shape by removing its XML element from the slide's shape tree.
    
    Args:
        slide: The slide object containing the shape
        shape: The shape object to remove
    """
    try:
        parent = shape.element.getparent()
        if parent is not None:
            parent.remove(shape.element)
    except Exception as e:
        print(f"Error removing shape: {e}")
        # Continue execution even if shape removal fails


def create_new_presentation_from_slides(slides_to_include):
    """
    Create a new presentation containing only the specified slides.
    
    Args:
        slides_to_include: List of slide objects to include in the new presentation
        
    Returns:
        A new Presentation object with only the specified slides
    """
    if not slides_to_include:
        return None
    
    # Create new presentation
    new_prs = Presentation()
    
    # Copy master slides and layouts from first presentation
    # Note: This is simplified - complete implementation would need to handle masters and layouts
    
    # Add each slide to the new presentation
    for slide in slides_to_include:
        # Get layout from original slide
        layout = slide.slide_layout
        
        # Add new slide with same layout
        new_slide = new_prs.slides.add_slide(layout)
        
        # Copy all shapes from source to destination
        for shape in slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_prs