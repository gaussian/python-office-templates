from __future__ import annotations

from io import BytesIO
from urllib.request import urlopen

from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl.drawing.image import Image as XLImage

from .exceptions import ImageError


def extract_image_url(text: str | None) -> str | None:
    """Return the image URL if *text* starts with the %image% directive."""
    if not text:
        return None
    stripped = text.strip()
    if not stripped.lower().startswith("%image%"):
        return None
    url = stripped[len("%image%"):].strip()
    return url or None


def should_replace_shape_with_image(shape) -> bool:
    """Return True if the shape text indicates an image placeholder."""
    if not hasattr(shape, "text_frame"):
        return False
    return extract_image_url(shape.text_frame.text) is not None


def should_replace_cell_with_image(cell) -> bool:
    """Return True if the cell value indicates an image placeholder."""
    if not isinstance(cell.value, str):
        return False
    return extract_image_url(cell.value) is not None

def replace_shape_with_image(shape, slide, url: str | None = None):
    """Replace *shape* with an image, keeping its size and position."""

    if url is None:
        if not hasattr(shape, "text_frame"):
            return
        url = extract_image_url(shape.text_frame.text)
    if not url:
        return

    try:
        with urlopen(url) as resp:
            data = resp.read()
    except Exception as e:  # pragma: no cover - network issues
        raise ImageError(f"Failed to download image from {url}: {e}")

    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    rotation = getattr(shape, "rotation", 0)

    pic = slide.shapes.add_picture(BytesIO(data), left, top, width=width, height=height)
    pic.rotation = rotation

    # Remove the original shape
    sp_tree = shape._element.getparent()
    sp_tree.remove(shape._element)

    return pic


def replace_cell_with_image(cell, worksheet, url: str | None = None):
    """Replace the cell's value with an image anchored at the cell."""

    if url is None:
        url = extract_image_url(cell.value if isinstance(cell.value, str) else None)
    if not url:
        return

    try:
        with urlopen(url) as resp:
            data = resp.read()
    except Exception as e:  # pragma: no cover - network issues
        raise ImageError(f"Failed to download image from {url}: {e}")

    img = XLImage(BytesIO(data))
    img.anchor = cell.coordinate
    worksheet.add_image(img)
    cell.value = None
    return img
