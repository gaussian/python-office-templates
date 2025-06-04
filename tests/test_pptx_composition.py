import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from template_reports.office_renderer import compose_pptx


# Dummy objects for testing
class DummyUser:
    def __init__(self, name, email):
        self.name = name
        self.email = email
        self._meta = True

    def __str__(self):
        return self.name


class TestPPTXComposition(unittest.TestCase):
    def setUp(self):
        # Create test presentations with layouts and slides
        self.temp_files = []
        self.template1_path = self._create_test_template1()
        self.template2_path = self._create_test_template2()
        
        # Test context
        self.context = {
            'user': DummyUser('Alice', 'alice@example.com'),
            'title': 'Test Presentation',
            'content': 'This is test content',
        }

    def tearDown(self):
        # Clean up temporary files
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_test_template1(self):
        """Create a test template with multiple layouts."""
        prs = Presentation()
        
        # Add a slide with title layout (for title layouts test)
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide1.shapes.title.text = "Title Layout"
        
        # Add a slide with content (for tagged layouts test)
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide2.shapes.title.text = "Content Slide"
        text_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        text_box.text_frame.text = "% layout content %"
        
        # Save to temp file
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def _create_test_template2(self):
        """Create another test template with different content."""
        prs = Presentation()
        
        # Add a slide for tagged layout
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        text_box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        text_box.text_frame.text = "% layout special %"
        
        # Add another slide with title for title layouts
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide2.shapes.title.text = "Special Layout"
        
        # Save to temp file
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_compose_with_master_layouts(self):
        """Test basic composition using master slide layouts."""
        slides = [
            {'layout': 'Title Slide', 'title': 'My Title'},
            {'layout': 'Title and Content', 'content': 'Some content'},
        ]
        
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slides=slides,
            global_context=self.context,
            output=output_file
        )
        
        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        
        # Check output file exists and has correct number of slides
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_with_tagged_layouts(self):
        """Test composition using tagged layouts."""
        slides = [
            {'layout': 'content', 'title': 'Content Title'},
            {'layout': 'special', 'content': 'Special content'},
        ]
        
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slides=slides,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True
        )
        
        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        
        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_with_title_layouts(self):
        """Test composition using slide titles as layout IDs."""
        slides = [
            {'layout': 'Title Layout', 'content': 'From title layout'},
            {'layout': 'Special Layout', 'content': 'From special layout'},
        ]
        
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slides=slides,
            global_context=self.context,
            output=output_file,
            use_all_slides_as_layouts_by_title=True
        )
        
        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        
        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_with_placeholders(self):
        """Test composition with placeholder processing."""
        slides = [
            {
                'layout': 'Title Slide',
                'placeholders': ['{{ title }}', '{{ user.name }}'],
                'title': 'Placeholder Test'
            },
        ]
        
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slides=slides,
            global_context=self.context,
            output=output_file
        )
        
        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_compose_with_missing_layout(self):
        """Test error handling when layout is not found."""
        slides = [
            {'layout': 'NonExistentLayout', 'content': 'Test'},
        ]
        
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slides=slides,
            global_context=self.context,
            output=output_file
        )
        
        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("Layout 'NonExistentLayout' not found" in error for error in errors))

    def test_compose_with_no_template_files(self):
        """Test error handling when no template files are provided."""
        slides = [{'layout': 'test', 'content': 'Test'}]
        
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[],
            slides=slides,
            global_context=self.context,
            output=output_file
        )
        
        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("No template files provided" in error for error in errors))

    def test_compose_with_no_slides(self):
        """Test error handling when no slides are specified."""
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slides=[],
            global_context=self.context,
            output=output_file
        )
        
        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("No slides specified" in error for error in errors))


if __name__ == '__main__':
    unittest.main()