import os
import tempfile
import unittest
from pptx import Presentation
from pptx.util import Inches
from template_reports.pptx_renderer.renderer import render_pptx
from template_reports.pptx_renderer.exceptions import (
    UnterminatedTagException,
    PermissionDeniedException,
    UnresolvedTagError,
)
import datetime


# Dummy objects for testing renderer behavior.
class DummyUser:
    def __init__(self, name, email):
        self.name = name
        self.email = email
        self._meta = True  # Simulate a Django model

    def __str__(self):
        return self.name


class DummyRequestUser:
    def has_perm(self, perm, obj):
        # Deny permission if object's name contains "deny"
        if hasattr(obj, "name") and "deny" in obj.name.lower():
            return False
        return True


class TestRendererUnit(unittest.TestCase):
    def setUp(self):
        # Create a minimal PPTX file with one slide and one text box.
        self.prs = Presentation()
        slide_layout = self.prs.slide_layouts[5]  # use a blank layout
        self.slide = self.prs.slides.add_slide(slide_layout)
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1)
        )
        self.textframe = self.textbox.text_frame
        # Insert a paragraph with a pure placeholder.
        p = self.textframe.paragraphs[0]
        p.text = "Hello, {{ user.name }}"

        # Create a temporary file to save the PPTX.
        self.temp_input = tempfile.mktemp(suffix=".pptx")
        self.temp_output = tempfile.mktemp(suffix=".pptx")
        self.prs.save(self.temp_input)

        # Prepare dummy context.
        self.user = DummyUser("Alice", "alice@example.com")
        self.context = {
            "user": self.user,
            "now": datetime.datetime(2025, 2, 18, 12, 0, 0),
        }
        self.request_user = DummyRequestUser()

    def tearDown(self):
        # Clean up temporary files.
        if os.path.exists(self.temp_input):
            os.remove(self.temp_input)
        if os.path.exists(self.temp_output):
            os.remove(self.temp_output)

    def test_textbox_pure_placeholder_normal_mode(self):
        # Test that pure placeholder in a text box is rendered correctly.
        rendered = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            request_user=self.request_user,
            check_permissions=False,
        )
        # Load output and check text from slide[0], shape[0].
        prs_out = Presentation(rendered)
        shape = prs_out.slides[0].shapes[0]
        txt = shape.text_frame.text
        self.assertEqual(txt, "Hello, Alice")

    def test_textbox_mixed_text_with_date(self):
        # Create a PPTX with mixed text containing a date placeholder.
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text_frame.text = "Date: {{ now | MMMM dd, YYYY }}"
        temp_in = tempfile.mktemp(suffix=".pptx")
        temp_out = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_in)

        context = {"now": datetime.datetime(2025, 2, 18, 12, 0, 0)}
        rendered = render_pptx(
            temp_in,
            context,
            temp_out,
            request_user=self.request_user,
            check_permissions=False,
        )
        prs_out = Presentation(rendered)
        shape = prs_out.slides[0].shapes[0]
        text = shape.text_frame.text
        expected = "Date: " + datetime.datetime(2025, 2, 18, 12, 0, 0).strftime(
            "%B %d, %Y"
        )
        self.assertEqual(text, expected)
        os.remove(temp_in)
        os.remove(temp_out)

    def test_mixed_text_list_joining(self):
        # Test mixed text where placeholder resolves to a list.
        # For example: "Emails: {{ user.email }}"
        # We'll set user.email to be a list for this test.
        self.user.email = ["alice@example.com", "alice.alt@example.com"]
        self.textframe.paragraphs[0].text = "Emails: {{ user.email }}"
        self.prs.save(self.temp_input)
        rendered = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            request_user=self.request_user,
            check_permissions=False,
        )
        prs_out = Presentation(rendered)
        shape = prs_out.slides[0].shapes[0]
        text = shape.text_frame.text
        self.assertEqual(text, "Emails: alice@example.com, alice.alt@example.com")

    def test_permission_denied_pure(self):
        self.user.name = "DenyUser"
        # Retrieve entire user (Django-like object), forcing permission denial:
        self.textframe.paragraphs[0].text = "{{ user }}"
        self.prs.save(self.temp_input)
        with self.assertRaises(Exception) as cm:
            render_pptx(
                self.temp_input,
                self.context,
                self.temp_output,
                request_user=self.request_user,
                check_permissions=True,
            )
        self.assertIn("Permission denied", str(cm.exception))


if __name__ == "__main__":
    unittest.main()
