import unittest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from template_reports.pptx_renderer.charts import process_chart


class DummyPlot:
    def __init__(self, categories):
        # categories is a list of strings that may contain placeholders.
        self.categories = categories


class DummySeries:
    def __init__(self, name, values):
        # name is a string (can contain placeholders); values is a list of floats.
        self.name = name
        self.values = values


class DummyChart:
    def __init__(self, plots, series):
        self.plots = plots  # List of DummyPlot objects.
        self.series = series  # List of DummySeries objects.
        self.replaced_data = (
            None  # Will hold the ChartData after replace_data() is called.
        )

    def replace_data(self, chart_data):
        # Capture the new ChartData.
        self.replaced_data = chart_data


class TestProcessChart(unittest.TestCase):
    def test_chart_data_replacement(self):
        # Define a simple context.
        context = {"test": "Replaced"}
        # Create a dummy plot with placeholders.
        # e.g., "Category {{ test }}" should become "Category Replaced"
        dummy_plot = DummyPlot(categories=["Category {{ test }}", "Static"])
        # Create a dummy series with a name containing a placeholder.
        dummy_series = DummySeries("Series {{ test }}", [1.0, 2.0])
        # Create the dummy chart with one plot and one series.
        dummy_chart = DummyChart(plots=[dummy_plot], series=[dummy_series])

        # Call process_chart. This should update our dummy_chart.replaced_data.
        process_chart(dummy_chart, context, perm_user=None)

        # Verify that process_chart replaced the placeholders.
        self.assertIsNotNone(dummy_chart.replaced_data)
        chart_data = dummy_chart.replaced_data

        # Convert categories into a list before comparing.
        expected_categories = ["Category Replaced", "Static"]
        actual_categories = list(c.label for c in chart_data.categories)
        self.assertEqual(actual_categories, expected_categories)

        # Check that series name is processed and values remain unchanged.
        self.assertEqual(len(chart_data._series), 1)
        series_obj = chart_data._series[0]
        self.assertEqual(series_obj.name, "Series Replaced")
        self.assertEqual(list(series_obj.values), [1.0, 2.0])


class TestProcessChartFull(unittest.TestCase):
    def test_chart_data_replacement(self):
        # Create a Presentation and add a chart with placeholder values.
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Build initial chart data with placeholders.
        chart_data = ChartData()
        chart_data.categories = ["Category {{ test }}", "Static"]
        chart_data.add_series("Series {{ test }}", (1.0, 2.0))

        # Add a chart shape (e.g. a clustered column chart).
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
            chart_data,
        )
        chart = chart_shape.chart

        # Process the chart with a context where "{{ test }}" becomes "Replaced".
        context = {"test": "Replaced"}
        process_chart(chart, context, perm_user=None)

        # Now check that the chart's data was updated:
        # Retrieve the categories from the first plot.
        new_categories = [str(cat) for cat in chart.plots[0].categories]
        self.assertEqual(new_categories, ["Category Replaced", "Static"])

        # Check series: the series name should be updated and values should remain the same.
        self.assertEqual(len(chart.series), 1)
        series_obj = chart.series[0]
        self.assertEqual(series_obj.name, "Series Replaced")
        self.assertEqual(list(series_obj.values), [1.0, 2.0])


if __name__ == "__main__":
    unittest.main()
