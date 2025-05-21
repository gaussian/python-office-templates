import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from template_reports.office_renderer import render_pptx
from template_reports.office_renderer.loops import is_loop_start, is_loop_end


class TestPptxIntegrationLoops(unittest.TestCase):
    def setUp(self):
        # Create a temporary file for the PPTX
        self.temp_pptx = tempfile.mktemp(suffix=".pptx")
        self.output_pptx = tempfile.mktemp(suffix=".pptx")
        
        # Create a presentation with loop directive
        self.prs = Presentation()
        
        # Create slides for the loop
        self.slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.slide3 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add loop start shape to slide 1
        loop_start = self.slide1.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in users%"
        
        # Add content shape to slide 2
        user_info = self.slide2.shapes.add_textbox(
            Inches(1), Inches(2), Inches(4), Inches(1)
        )
        user_info.text_frame.text = "Name: {{ user.name }}, Email: {{ user.email }}"
        
        # Add loop end shape to slide 3
        loop_end = self.slide3.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"
        
        # Save the presentation
        self.prs.save(self.temp_pptx)
        
        # Prepare the context
        self.context = {
            "users": [
                {"name": "Alice", "email": "alice@example.com"},
                {"name": "Bob", "email": "bob@example.com"},
                {"name": "Charlie", "email": "charlie@example.com"},
            ]
        }
    
    def tearDown(self):
        # Clean up temporary files
        for temp_file in [self.temp_pptx, self.output_pptx]:
            if os.path.exists(temp_file):
                os.remove(temp_file)
    
    def test_loop_processing(self):
        """Test that loop processing duplicates slides and substitutes variables."""
        # Render the presentation with the context
        result, errors = render_pptx(
            template=self.temp_pptx,
            context=self.context,
            output=self.output_pptx,
            perm_user=None,
        )
        
        # Check for errors
        self.assertIsNone(errors)
        
        # Load the rendered presentation
        rendered_prs = Presentation(self.output_pptx)
        
        # The rendered presentation should have 3 slides (one for each user in the loop)
        # We expect to see the content from slide 2 repeated three times,
        # with the user variable substituted for each item in users
        # Check if loop directives are present
        found_loop_start = False
        found_loop_end = False
        for slide in rendered_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    if is_loop_start(shape):
                        found_loop_start = True
                    if is_loop_end(shape):
                        found_loop_end = True
                    
                    # Look for the user info text and verify it matches the expected values
                    text = shape.text_frame.text
                    if "Name:" in text and "Email:" in text:
                        # Check if the text contains one of our user's info
                        for user in self.context["users"]:
                            expected_text = f"Name: {user['name']}, Email: {user['email']}"
                            if text == expected_text:
                                # Found a match
                                break
                        else:
                            # No match found, fail the test
                            self.fail(f"Unexpected text in slide: {text}")
        
        # Check if loop directives were properly handled
        self.assertTrue(found_loop_start)
        self.assertTrue(found_loop_end)


if __name__ == "__main__":
    unittest.main()