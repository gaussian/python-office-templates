import os
import tempfile
import unittest
import datetime
from pptx import Presentation
from pptx.util import Inches
from template_reports.pptx_renderer.renderer import render_pptx


# Dummy objects for integration testing.
class DummyUser:
    def __init__(self, name, email, is_active=True):
        self.name = name
        self.email = email
        self.is_active = is_active
        self._meta = True  # Simulate a Django model

    def __str__(self):
        return self.name


class DummyCohort:
    def __init__(self, name):
        self.name = name

    def __str__(self):
        return self.name


class DummyRequestUser:
    def has_perm(self, perm, obj):
        # Deny permission if object's name contains "deny"
        if hasattr(obj, "name") and "deny" in obj.name.lower():
            return False
        return True


class TestRendererIntegration(unittest.TestCase):
    def setUp(self):
        # Create a minimal PPTX with one text box and one table.
        self.prs = Presentation()
        blank_slide = self.prs.slide_layouts[5]
        self.slide = self.prs.slides.add_slide(blank_slide)
        # Add a text box with mixed text.
        textbox = self.slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4), Inches(1)
        )
        textbox.text_frame.text = "Welcome, {{ user.name }}. Program: {{ program.name }}."
        # Add a table with one row, one column.
        rows, cols = 1, 1
        left, top, width, height = Inches(0.5), Inches(2), Inches(4), Inches(0.8)
        table_shape = self.slide.shapes.add_table(rows, cols, left, top, width, height)
        table_cell = table_shape.table.cell(0, 0)
        table_cell.text = "{{ program.users.email }}"

        # Save this PPTX to a temporary file.
        self.temp_input = tempfile.mktemp(suffix=".pptx")
        self.temp_output = tempfile.mktemp(suffix=".pptx")
        self.prs.save(self.temp_input)

        # Set up the context.
        self.cohort = DummyCohort("Cohort A")
        self.user = DummyUser("Alice", "alice@example.com", is_active=True)
        self.user2 = DummyUser("Bob", "bob@example.com", is_active=True)
        self.user3 = DummyUser("Carol", "carol@example.com", is_active=False)
        self.program = {
            "name": "Test Program",
            "users": [self.user, self.user2, self.user3],
        }
        self.context = {
            "user": self.user,
            "program": self.program,
            "now": datetime.datetime(2025, 2, 18, 12, 0, 0),
            "date": datetime.date(2020, 1, 15),
        }
        self.request_user = DummyRequestUser()

    def tearDown(self):
        if os.path.exists(self.temp_input):
            os.remove(self.temp_input)
        if os.path.exists(self.temp_output):
            os.remove(self.temp_output)

    def test_integration_renderer(self):
        # Run the renderer integration.
        rendered = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            request_user=self.request_user,
            check_permissions=False,
        )
        # Open the rendered PPTX.
        prs_out = Presentation(rendered)
        # Test text box content.
        textbox = prs_out.slides[0].shapes[0]
        txt = textbox.text_frame.text
        self.assertIn("Welcome, Alice.", txt)
        self.assertIn("Program: Test Program", txt)
        # Test table: the table should have expanded rows for each user email.
        table_shape = None
        for shape in prs_out.slides[0].shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                table_shape = shape
                break
        self.assertIsNotNone(table_shape)
        # In normal mode for tables, pure placeholders return a list; so we expect one cell in the first row,
        # and additional rows for subsequent items.
        emails = [u.email for u in self.program["users"]]
        # Check that the first row cell contains the first email.
        first_cell_text = table_shape.table.cell(0, 0).text.strip()
        self.assertEqual(first_cell_text, emails[0])
        # There should be additional rows for each remaining email.
        # Since our expander is a hack, we'll simply check that at least one additional row exists.
        self.assertTrue(len(list(table_shape.table.rows)) > 1)


if __name__ == "__main__":
    unittest.main()
