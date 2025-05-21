import unittest
from unittest.mock import MagicMock, patch

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from template_reports.office_renderer.loops import (
    extract_loop_directive,
    is_loop_end,
    is_loop_start,
)
from template_reports.office_renderer.pptx import process_loops


class TestLoopDirectives(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(2)
        )

    def test_extract_loop_directive(self):
        """Test the extraction of loop variable and collection from directive."""
        # Valid loop directive
        self.textbox.text_frame.text = "%loop user in users%"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "users")

        # Invalid loop directive
        self.textbox.text_frame.text = "Not a loop directive"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertIsNone(variable)
        self.assertIsNone(collection)

        # Empty text
        variable, collection = extract_loop_directive(None)
        self.assertIsNone(variable)
        self.assertIsNone(collection)

    def test_is_loop_start(self):
        """Test the detection of loop start directive."""
        # Valid loop start
        self.textbox.text_frame.text = "%loop user in users%"
        self.assertTrue(is_loop_start(self.textbox))

        # Not a loop start
        self.textbox.text_frame.text = "Not a loop start"
        self.assertFalse(is_loop_start(self.textbox))

        # Shape has no text_frame
        shape = MagicMock()
        delattr(shape, 'text_frame')
        self.assertFalse(is_loop_start(shape))

    def test_is_loop_end(self):
        """Test the detection of loop end directive."""
        # Valid loop end
        self.textbox.text_frame.text = "%endloop%"
        self.assertTrue(is_loop_end(self.textbox))

        # Not a loop end
        self.textbox.text_frame.text = "Not a loop end"
        self.assertFalse(is_loop_end(self.textbox))

        # Shape has no text_frame
        shape = MagicMock()
        delattr(shape, 'text_frame')
        self.assertFalse(is_loop_end(shape))


class TestLoopProcessing(unittest.TestCase):
    def setUp(self):
        self.context = {"users": [
            {"name": "Alice", "email": "alice@example.com"},
            {"name": "Bob", "email": "bob@example.com"}
        ]}
        
    @patch("template_reports.office_renderer.pptx.process_loops")
    @patch("template_reports.office_renderer.pptx.Presentation")
    def test_process_loops_called(self, mock_presentation, mock_process_loops):
        """Test that process_loops is called from render_pptx."""
        from template_reports.office_renderer.pptx import render_pptx

        # Prepare mock return value for process_loops
        mock_process_loops.return_value = []
        
        # Setup mock for Presentation
        mock_prs = MagicMock()
        mock_presentation.return_value = mock_prs
        
        # Call render_pptx with a string template path
        render_pptx("template.pptx", {}, "output.pptx", None)
        
        # Verify process_loops was called
        mock_process_loops.assert_called_once()
    
    def test_extract_loop_directive_integration(self):
        """Test extract_loop_directive directly."""
        # Valid loop
        directive = "%loop user in users%"
        variable, collection = extract_loop_directive(directive)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "users")
        
        # Invalid format
        directive = "%loop invalid directive%"
        variable, collection = extract_loop_directive(directive)
        self.assertIsNone(variable)
        self.assertIsNone(collection)
        
        # Empty
        directive = ""
        variable, collection = extract_loop_directive(directive)
        self.assertIsNone(variable)
        self.assertIsNone(collection)


if __name__ == "__main__":
    unittest.main()