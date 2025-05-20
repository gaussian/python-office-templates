import os
import tempfile
import unittest
from unittest.mock import patch

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from PIL import Image

from template_reports.office_renderer.images import replace_shape_with_image, ImageError


class TestSlideDownloadedImage(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(2)
        )
        # Create a small red image on disk
        img = Image.new("RGB", (4, 4), color="red")
        self.temp_image = tempfile.mktemp(suffix=".png")
        img.save(self.temp_image)
        self.textbox.text_frame.text = f"%image% file://{self.temp_image}"

    def tearDown(self):
        if os.path.exists(self.temp_image):
            os.remove(self.temp_image)

    def test_replace_shape_with_image(self):
        original_width = self.textbox.width
        original_height = self.textbox.height

        replace_shape_with_image(self.textbox, self.slide)

        # Only one non-placeholder shape should remain and it should be a picture
        shapes = [
            shape for shape in self.slide.shapes if not shape.is_placeholder
        ]
        self.assertEqual(len(shapes), 1)
        pic = shapes[0]
        self.assertEqual(pic.shape_type, MSO_SHAPE_TYPE.PICTURE)
        self.assertEqual(pic.width, original_width)
        self.assertEqual(pic.height, original_height)

    @patch("template_reports.office_renderer.images.urlopen")
    def test_invalid_url_raises_image_error(self, mock_urlopen):
        mock_urlopen.side_effect = Exception("boom")
        self.textbox.text_frame.text = "%image% http://example.com/foo.png"
        with self.assertRaises(ImageError):
            replace_shape_with_image(self.textbox, self.slide)


if __name__ == "__main__":
    unittest.main()