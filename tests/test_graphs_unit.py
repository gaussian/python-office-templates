import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import compose_graphs_pptx


class TestGraphsUnit(unittest.TestCase):
    """Unit tests for graph composition functionality."""

    def setUp(self):
        self.temp_files = []
        self.template_path = self._create_test_template()
        self.context = {
            "company": "Test Corp",
            "node_prefix": "Node",
        }

    def tearDown(self):
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_test_template(self):
        """Create a simple test template."""
        prs = Presentation()
        
        # Add a blank slide as a graph layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add layout tag
        layout_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(0.5)
        )
        layout_box.text_frame.text = "% layout graph %"

        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_simple_graph_with_nodes_only(self):
        """Test creating a simple graph with just nodes, no edges."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "node1",
                        "name": "First Node",
                        "detail": "Details about first node",
                        "position": {"x": 1, "y": 1},
                    },
                    {
                        "id": "node2",
                        "name": "Second Node",
                        "detail": "Details about second node",
                        "position": {"x": 4, "y": 1},
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        # Verify the output
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)
        
        # Check that slide has shapes (nodes)
        slide = prs.slides[0]
        self.assertGreater(len(slide.shapes), 0)

    def test_graph_with_nodes_and_edges(self):
        """Test creating a graph with nodes and edges."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "A",
                        "name": "Node A",
                        "position": {"x": 1, "y": 2},
                    },
                    {
                        "id": "B",
                        "name": "Node B",
                        "position": {"x": 4, "y": 2},
                    },
                    {
                        "id": "C",
                        "name": "Node C",
                        "position": {"x": 7, "y": 2},
                    },
                ],
                "edges": [
                    {"from": "A", "to": "B", "label": "connects to"},
                    {"from": "B", "to": "C", "label": "flows to"},
                ],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

    def test_multiple_graphs_create_multiple_slides(self):
        """Test that multiple graphs create multiple slides."""
        graphs = [
            {
                "nodes": [
                    {"id": "1", "name": "Node 1", "position": {"x": 1, "y": 1}},
                ],
                "edges": [],
            },
            {
                "nodes": [
                    {"id": "2", "name": "Node 2", "position": {"x": 1, "y": 1}},
                ],
                "edges": [],
            },
            {
                "nodes": [
                    {"id": "3", "name": "Node 3", "position": {"x": 1, "y": 1}},
                ],
                "edges": [],
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 3)

    def test_template_variable_processing_in_nodes(self):
        """Test that template variables in node names are processed."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "test",
                        "name": "{{ company }}",
                        "detail": "Prefix: {{ node_prefix }}",
                        "position": {"x": 1, "y": 1},
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check that variables were processed
        prs = Presentation(output_file)
        slide = prs.slides[0]
        
        # Find text containing processed variable
        found_company = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text
                if "Test Corp" in text:
                    found_company = True
                    break
                    
        self.assertTrue(found_company, "Template variable should be processed")

    def test_edge_label_processing(self):
        """Test that edge labels with template variables are processed."""
        graphs = [
            {
                "nodes": [
                    {"id": "A", "name": "Node A", "position": {"x": 1, "y": 1}},
                    {"id": "B", "name": "Node B", "position": {"x": 4, "y": 1}},
                ],
                "edges": [
                    {"from": "A", "to": "B", "label": "{{ company }} edge"},
                ],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_node_without_position_error(self):
        """Test error handling when node is missing position."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "bad_node",
                        "name": "Bad Node",
                        # Missing position
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'position'" in error for error in errors))

    def test_node_without_id_error(self):
        """Test error handling when node is missing id."""
        graphs = [
            {
                "nodes": [
                    {
                        # Missing id
                        "name": "Node",
                        "position": {"x": 1, "y": 1},
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'id'" in error for error in errors))

    def test_node_without_name_error(self):
        """Test error handling when node is missing name."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "test",
                        # Missing name
                        "position": {"x": 1, "y": 1},
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'name'" in error for error in errors))

    def test_edge_with_unknown_node_error(self):
        """Test error handling when edge references unknown node."""
        graphs = [
            {
                "nodes": [
                    {"id": "A", "name": "Node A", "position": {"x": 1, "y": 1}},
                ],
                "edges": [
                    {"from": "A", "to": "B"},  # B doesn't exist
                ],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("unknown target node" in error for error in errors))

    def test_graph_missing_nodes_key_error(self):
        """Test error handling when graph is missing nodes key."""
        graphs = [
            {
                # Missing nodes
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("Missing 'nodes' key" in error for error in errors))

    def test_graph_missing_edges_key_error(self):
        """Test error handling when graph is missing edges key."""
        graphs = [
            {
                "nodes": [
                    {"id": "test", "name": "Node", "position": {"x": 1, "y": 1}},
                ],
                # Missing edges
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("Missing 'edges' key" in error for error in errors))

    def test_no_graphs_error(self):
        """Test error handling when no graphs are provided."""
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=[],
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("No graphs specified" in error for error in errors))

    def test_no_template_files_error(self):
        """Test error handling when no template files are provided."""
        graphs = [
            {
                "nodes": [
                    {"id": "test", "name": "Node", "position": {"x": 1, "y": 1}},
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("No template files provided" in error for error in errors))

    def test_node_with_parent_placeholder(self):
        """Test that nodes with parent key are accepted (placeholder functionality)."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "parent",
                        "name": "Parent Node",
                        "position": {"x": 1, "y": 1},
                    },
                    {
                        "id": "child",
                        "name": "Child Node",
                        "position": {"x": 2, "y": 2},
                        "parent": "parent",  # This is placeholder for now
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        # Should succeed - parent is accepted but not yet fully implemented
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_node_without_detail(self):
        """Test that nodes without detail field work correctly."""
        graphs = [
            {
                "nodes": [
                    {
                        "id": "simple",
                        "name": "Simple Node",
                        # No detail field
                        "position": {"x": 1, "y": 1},
                    },
                ],
                "edges": [],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_edge_without_label(self):
        """Test that edges without label field work correctly."""
        graphs = [
            {
                "nodes": [
                    {"id": "A", "name": "Node A", "position": {"x": 1, "y": 1}},
                    {"id": "B", "name": "Node B", "position": {"x": 4, "y": 1}},
                ],
                "edges": [
                    {"from": "A", "to": "B"},  # No label
                ],
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_graphs_pptx(
            template_files=[self.template_path],
            graphs=graphs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)


if __name__ == "__main__":
    unittest.main()
