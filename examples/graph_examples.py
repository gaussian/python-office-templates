"""
Example demonstrating how to create node/edge graphs in PowerPoint presentations.

This example shows various types of graphs that can be created using compose_pptx:
- Software architecture diagram
- Workflow/process diagram
- Network topology
- Organizational chart
"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from office_templates.office_renderer import compose_pptx


def create_template():
    """Create a simple graph template."""
    prs = Presentation()
    
    # Create a blank slide with a graph layout tag
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add layout tag
    layout_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(4), Inches(0.5)
    )
    layout_box.text_frame.text = "% layout graph %"
    
    # Add title placeholder
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0), Inches(8), Inches(0.5)
    )
    title_box.text_frame.text = "Graph: {{ graph_title }}"
    
    # Save template
    template_path = tempfile.mktemp(suffix=".pptx")
    prs.save(template_path)
    return template_path


def example_software_architecture():
    """Example: Software architecture diagram."""
    return {
        "nodes": [
            {
                "id": "user",
                "name": "User",
                "detail": "Web Browser",
                "position": {"x": 1, "y": 2},
            },
            {
                "id": "cdn",
                "name": "CDN",
                "detail": "CloudFront",
                "position": {"x": 3, "y": 2},
            },
            {
                "id": "lb",
                "name": "Load Balancer",
                "detail": "AWS ALB",
                "position": {"x": 5, "y": 2},
            },
            {
                "id": "app",
                "name": "Application",
                "detail": "Node.js",
                "position": {"x": 7, "y": 1.5},
            },
            {
                "id": "cache",
                "name": "Cache",
                "detail": "Redis",
                "position": {"x": 9, "y": 1},
            },
            {
                "id": "db",
                "name": "Database",
                "detail": "PostgreSQL",
                "position": {"x": 9, "y": 2.5},
            },
        ],
        "edges": [
            {"from": "user", "to": "cdn", "label": "HTTPS"},
            {"from": "cdn", "to": "lb", "label": "Forward"},
            {"from": "lb", "to": "app", "label": "Route"},
            {"from": "app", "to": "cache", "label": "Check"},
            {"from": "app", "to": "db", "label": "Query"},
        ],
    }


def example_workflow():
    """Example: Workflow/process diagram."""
    return {
        "nodes": [
            {
                "id": "start",
                "name": "Start",
                "detail": "Order Received",
                "position": {"x": 1, "y": 2},
            },
            {
                "id": "validate",
                "name": "Validate",
                "detail": "Check Inventory",
                "position": {"x": 3, "y": 2},
            },
            {
                "id": "authorize",
                "name": "Authorize Payment",
                "detail": "Credit Card",
                "position": {"x": 5, "y": 2},
            },
            {
                "id": "fulfill",
                "name": "Fulfill Order",
                "detail": "Ship Products",
                "position": {"x": 7, "y": 2},
            },
            {
                "id": "complete",
                "name": "Complete",
                "detail": "Order Shipped",
                "position": {"x": 9, "y": 2},
            },
        ],
        "edges": [
            {"from": "start", "to": "validate", "label": "Submit"},
            {"from": "validate", "to": "authorize", "label": "Valid"},
            {"from": "authorize", "to": "fulfill", "label": "Approved"},
            {"from": "fulfill", "to": "complete", "label": "Done"},
        ],
    }


def example_network_topology():
    """Example: Network topology diagram."""
    return {
        "nodes": [
            {"id": "internet", "name": "Internet", "position": {"x": 4, "y": 0.5}},
            {"id": "router", "name": "Router", "detail": "Main Gateway", "position": {"x": 4, "y": 2}},
            {"id": "switch1", "name": "Switch 1", "detail": "Floor 1", "position": {"x": 2, "y": 3.5}},
            {"id": "switch2", "name": "Switch 2", "detail": "Floor 2", "position": {"x": 6, "y": 3.5}},
            {"id": "server1", "name": "Web Server", "position": {"x": 1, "y": 5}},
            {"id": "server2", "name": "DB Server", "position": {"x": 3, "y": 5}},
            {"id": "server3", "name": "App Server", "position": {"x": 5, "y": 5}},
            {"id": "server4", "name": "File Server", "position": {"x": 7, "y": 5}},
        ],
        "edges": [
            {"from": "internet", "to": "router", "label": "WAN"},
            {"from": "router", "to": "switch1", "label": "1Gbps"},
            {"from": "router", "to": "switch2", "label": "1Gbps"},
            {"from": "switch1", "to": "server1", "label": "100Mbps"},
            {"from": "switch1", "to": "server2", "label": "100Mbps"},
            {"from": "switch2", "to": "server3", "label": "100Mbps"},
            {"from": "switch2", "to": "server4", "label": "100Mbps"},
        ],
    }


def example_org_chart():
    """Example: Organizational chart."""
    return {
        "nodes": [
            {
                "id": "ceo",
                "name": "CEO",
                "detail": "Chief Executive",
                "position": {"x": 4, "y": 1},
            },
            {
                "id": "cto",
                "name": "CTO",
                "detail": "Technology",
                "position": {"x": 2, "y": 3},
            },
            {
                "id": "cfo",
                "name": "CFO",
                "detail": "Finance",
                "position": {"x": 6, "y": 3},
            },
            {
                "id": "eng",
                "name": "Engineering Manager",
                "position": {"x": 1, "y": 5},
            },
            {
                "id": "devops",
                "name": "DevOps Manager",
                "position": {"x": 3, "y": 5},
            },
        ],
        "edges": [
            {"from": "ceo", "to": "cto", "label": "Reports"},
            {"from": "ceo", "to": "cfo", "label": "Reports"},
            {"from": "cto", "to": "eng", "label": "Manages"},
            {"from": "cto", "to": "devops", "label": "Manages"},
        ],
    }


def main():
    """Generate example graph presentations."""
    print("Creating graph examples...")
    
    # Create template
    template_path = create_template()
    print(f"Created template at: {template_path}")
    
    # Define slide specs with different graph types
    slide_specs = [
        {
            "layout": "graph",
            "graph": example_software_architecture(),
        },
        {
            "layout": "graph",
            "graph": example_workflow(),
        },
        {
            "layout": "graph",
            "graph": example_network_topology(),
        },
        {
            "layout": "graph",
            "graph": example_org_chart(),
        },
    ]
    
    # Global context for template variables
    context = {
        "graph_title": "System Diagrams",
        "company": "Example Corp",
    }
    
    # Generate output
    output_path = tempfile.mktemp(suffix=".pptx")
    result, errors = compose_pptx(
        template_files=[template_path],
        slide_specs=slide_specs,
        global_context=context,
        output=output_path,
        use_tagged_layouts=True,
    )
    
    if errors:
        print("Errors occurred:")
        for error in errors:
            print(f"  - {error}")
    else:
        print(f"Successfully created presentation at: {output_path}")
        
        # Verify output
        prs = Presentation(output_path)
        print(f"Number of slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            print(f"  Slide {i+1}: {len(slide.shapes)} shapes")
            
        print("\nYou can open the file to view the generated graphs:")
        print(f"  {output_path}")
    
    # Clean up template
    if os.path.exists(template_path):
        os.remove(template_path)


if __name__ == "__main__":
    main()
